import streamlit as st
import google.generativeai as genai
from datetime import datetime
import speech_recognition as sr
import PyPDF2
from pptx import Presentation
import time 
import io

# --- PAGE CONFIG ---
st.set_page_config(
    page_title="Texer",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Hide Streamlit navigation and footer
hide_streamlit_style = """
<style>
#MainMenu {visibility: hidden;}
footer {visibility: hidden;}
header {visibility: hidden;}
</style>
"""
st.markdown(hide_streamlit_style, unsafe_allow_html=True)

# Configure Gemini API
GEMINI_API_KEY = "AIzaSyD6zblME2aAIq3PEs07LGHAaWzmLXoALpM"
genai.configure(api_key=GEMINI_API_KEY)

# --- CUSTOM CSS ---
custom_css = """
<style>
    @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@400;500;600;700&display=swap');
    
    * {
        font-family: 'Poppins', sans-serif !important;
    }
    
    [data-testid="stAppViewContainer"] {
        background: linear-gradient(135deg, #f5f7fa 0%, #c3cfe2 100%);
        padding: 0 !important;
    }
    
    .header-container {
        position: relative;
        margin-top: -50px;
        margin-bottom: -10px;
        text-align: center;
        animation: fadeIn 0.8s ease-in-out;
    }
    
    @keyframes fadeIn {
        from { opacity: 0; transform: translateY(-20px); }
        to { opacity: 1; transform: translateY(0); }
    }
    
    .summary-section {
        margin-top: -30px;
    }
    
    .summary-box:hover {
        transform: translateY(-3px);
        box-shadow: 0 12px 40px rgba(31, 38, 135, 0.25);
    }
    
    .file-summary {
        border-left: 4px solid #6e45e2;
        padding-left: 1rem;
        margin-bottom: 1.5rem;
    }
    
    .stButton>button {
        border-radius: 14px !important;
        padding: 12px 24px !important;
        font-weight: 600 !important;
        transition: all 0.3s ease !important;
        border: none !important;
        margin: 8px 0 !important;
        font-size: 0.95rem !important;
        box-shadow: 0 4px 12px rgba(0,0,0,0.1) !important;
    }
    
    .primary-btn {
        background: linear-gradient(90deg, #6e45e2 0%, #89d3ce 100%) !important;
        color: white !important;
    }
    
    .primary-btn:hover {
        background: linear-gradient(90deg, #5d35d1 0%, #79c3be 100%) !important;
        transform: translateY(-2px) scale(1.02) !important;
        box-shadow: 0 8px 24px rgba(110, 69, 226, 0.3) !important;
    }
    
    .secondary-btn {
        background: linear-gradient(90deg, #f093fb 0%, #f5576c 100%) !important;
        color: white !important;
    }
    
    .secondary-btn:hover {
        background: linear-gradient(90deg, #e083eb 0%, #e5475c 100%) !important;
        transform: translateY(-2px) scale(1.02) !important;
        box-shadow: 0 8px 24px rgba(240, 147, 251, 0.3) !important;
    }
    
    .stTextArea>div>div>textarea {
        border-radius: 16px !important;
        padding: 16px !important;
        box-shadow: 0 4px 20px rgba(0,0,0,0.08) !important;
        border: 1px solid rgba(110, 69, 226, 0.3) !important;
        min-height: 150px;
        background: rgba(255, 255, 255, 0.9) !important;
        transition: all 0.3s ease;
    }
    
    .stTextArea>div>div>textarea:focus {
        border: 1px solid #6e45e2 !important;
        box-shadow: 0 4px 24px rgba(110, 69, 226, 0.15) !important;
    }
    
    .stFileUploader>div>div {
        border-radius: 16px !important;
        border: 2px dashed #6e45e2 !important;
        background: rgba(110, 69, 226, 0.05) !important;
        transition: all 0.3s ease;
    }
    
    .stFileUploader>div>div:hover {
        background: rgba(110, 69, 226, 0.1) !important;
        box-shadow: 0 4px 20px rgba(110, 69, 226, 0.1) !important;
    }
    
    .word-count {
        font-size: 0.85rem;
        color: #6e45e2;
        font-weight: 600;
        margin-top: -10px;
        margin-bottom: 15px;
        text-align: right;
    }
    
    [data-testid="stSidebar"] {
        background: linear-gradient(135deg, rgba(245, 247, 250, 0.9) 0%, rgba(195, 207, 226, 0.9) 100%) !important;
        padding: 1.5rem !important;
        box-shadow: 5px 0 15px rgba(0,0,0,0.05);
        border-right: 1px solid rgba(255, 255, 255, 0.3);
    }
    
    .sidebar-section {
        margin-bottom: 1.5rem;
        animation: slideIn 0.5s ease-out;
    }
    
    @keyframes slideIn {
        from { opacity: 0; transform: translateX(-20px); }
        to { opacity: 1; transform: translateX(0); }
    }
    
    .sidebar-title {
        color: #6e45e2;
        font-size: 1.1rem;
        font-weight: 700;
        margin-bottom: 0.75rem;
    }
    
    .history-item {
        padding: 1rem;
        margin-bottom: 1rem;
        border-radius: 12px;
        background: rgba(255, 255, 255, 0.7);
        box-shadow: 0 4px 12px rgba(0,0,0,0.05);
        transition: all 0.3s ease;
    }
    
    .history-item:hover {
        transform: translateY(-2px);
        box-shadow: 0 8px 16px rgba(0,0,0,0.1);
    }
    
    .history-title {
        font-weight: 600;
        margin-bottom: 0.25rem;
        color: #333;
    }
    
    .history-date {
        font-size: 0.75rem;
        color: #666;
        margin-bottom: 0.5rem;
    }
    
    .button-container {
        display: flex;
        gap: 1rem;
        margin-bottom: 1.5rem;
    }
    
    .divider {
        height: 1px;
        background: linear-gradient(90deg, transparent, rgba(110, 69, 226, 0.3), transparent);
        margin: 1.5rem 0;
    }
    
    .file-meta {
        font-size: 0.9rem;
        color: #666;
        margin-bottom: 0.5rem;
    }
    
    .file-type-badge {
        display: inline-block;
        padding: 0.25rem 0.5rem;
        border-radius: 12px;
        font-size: 0.75rem;
        font-weight: 600;
        margin-right: 0.5rem;
    }
    
    .pdf-badge {
        background-color: #ff6b6b;
        color: white;
    }
    
    .pptx-badge {
        background-color: #4ecdc4;
        color: white;
    }
    
    .txt-badge {
        background-color: #ffbe76;
        color: white;
    }
    
    .processing-animation {
        display: flex;
        align-items: center;
        margin: 1rem 0;
    }
    
    .spinner {
        border: 4px solid rgba(110, 69, 226, 0.1);
        border-radius: 50%;
        border-top: 4px solid #6e45e2;
        width: 20px;
        height: 20px;
        animation: spin 1s linear infinite;
        margin-right: 10px;
    }
    
    @keyframes spin {
        0% { transform: rotate(0deg); }
        100% { transform: rotate(360deg); }
    }
    
    .processing-text {
        color: #6e45e2;
        font-weight: 500;
    }
</style>
"""
st.markdown(custom_css, unsafe_allow_html=True)

# --- INITIALIZE SESSION STATE ---
def init_session_state():
    default_values = {
        "summary_mode": "Concise",
        "length_factor": 0.5,
        "history": [],
        "input_text": "",
        "file_processed": False,
        "show_summary": False,
        "generated_summary": "",
        "file_summaries": [],
        "processing_files": False,
        "copied_idx": None
    }
    for key, value in default_values.items():
        if key not in st.session_state:
            st.session_state[key] = value

init_session_state()

# --- CORE FUNCTIONS ---
def get_summary_from_gemini(text, mode, length_factor):
    try:
        # Truncate very large text to avoid API limits
        max_length = 50000  # Gemini's token limit
        truncated_text = text[:max_length] if len(text) > max_length else text
        
        model = genai.GenerativeModel("gemini-1.5-pro-latest")
        prompt = f"Summarize the following text in a {mode.lower()} manner. " \
                 f"The summary should be approximately {int(length_factor * 100)}% of the original content length.\n\n" \
                 f"Important: Provide only the summary content, no additional commentary or labels.\n\n" \
                 f"{truncated_text}"
        response = model.generate_content(prompt)
        return response.text.strip() if response.text else "No summary generated"
    except Exception as e:
        return f"Error: {str(e)}"

def summarize_pdf(uploaded_file):
    try:
        reader = PyPDF2.PdfReader(io.BytesIO(uploaded_file.read()))
        num_pages = len(reader.pages)
        full_text = ""
        
        # Show processing status
        status_container = st.empty()
        progress_bar = st.progress(0)
        
        # Extract text from each page with progress
        for i, page in enumerate(reader.pages):
            status_container.markdown(f"""
            <div class="processing-animation">
                <div class="spinner"></div>
                <div class="processing-text">Processing page {i+1} of {num_pages}</div>
            </div>
            """, unsafe_allow_html=True)
            
            page_text = page.extract_text()
            if page_text:
                full_text += f"\n\n[Page {i+1}]\n{page_text}"
            progress_bar.progress((i + 1) / num_pages)
            time.sleep(0.1)  # Small delay for smoother progress
        
        status_container.empty()
        progress_bar.empty()
        
        # Generate summary
        with st.spinner(f"Generating summary for {uploaded_file.name}..."):
            summary = get_summary_from_gemini(
                full_text,
                st.session_state.summary_mode,
                st.session_state.length_factor
            )
        
        return {
            "type": "PDF",
            "title": uploaded_file.name,
            "summary": summary,
            "original_text": full_text[:1000] + ("..." if len(full_text) > 1000 else ""),
            "meta": {
                "Author": reader.metadata.author or "Unknown",
                "Pages": num_pages,
                "Created": reader.metadata.creation_date or "Unknown"
            }
        }
    except Exception as e:
        return {
            "type": "PDF",
            "title": uploaded_file.name,
            "summary": f"Error processing PDF: {str(e)}",
            "original_text": "",
            "meta": {}
        }

def summarize_pptx(uploaded_file):
    try:
        prs = Presentation(io.BytesIO(uploaded_file.read()))
        full_text = []
        num_slides = len(prs.slides)
        
        # Show processing status
        status_container = st.empty()
        progress_bar = st.progress(0)
        
        # Extract text from each slide with progress
        for i, slide in enumerate(prs.slides):
            status_container.markdown(f"""
            <div class="processing-animation">
                <div class="spinner"></div>
                <div class="processing-text">Processing slide {i+1} of {num_slides}</div>
            </div>
            """, unsafe_allow_html=True)
            
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                full_text.append(f"[Slide {i+1}]\n" + "\n".join(slide_text))
            
            progress_bar.progress((i + 1) / num_slides)
            time.sleep(0.1)
        
        status_container.empty()
        progress_bar.empty()
        
        extracted_text = "\n\n".join(full_text) if full_text else "No text found in presentation"
        
        # Generate summary
        with st.spinner(f"Generating summary for {uploaded_file.name}..."):
            summary = get_summary_from_gemini(
                extracted_text,
                st.session_state.summary_mode,
                st.session_state.length_factor
            )
        
        return {
            "type": "PPTX",
            "title": uploaded_file.name,
            "summary": summary,
            "original_text": extracted_text[:1000] + ("..." if len(extracted_text) > 1000 else ""),
            "meta": {
                "Slides": num_slides
            }
        }
    except Exception as e:
        return {
            "type": "PPTX",
            "title": uploaded_file.name,
            "summary": f"Error processing PPTX: {str(e)}",
            "original_text": "",
            "meta": {}
        }

def summarize_txt(uploaded_file):
    try:
        # Show processing status
        with st.spinner(f"Processing {uploaded_file.name}..."):
            extracted_text = uploaded_file.read().decode("utf-8")
            
            # Generate summary
            summary = get_summary_from_gemini(
                extracted_text,
                st.session_state.summary_mode,
                st.session_state.length_factor
            )
        
        return {
            "type": "TXT",
            "title": uploaded_file.name,
            "summary": summary,
            "original_text": extracted_text[:1000] + ("..." if len(extracted_text) > 1000 else ""),
            "meta": {
                "Size": f"{len(extracted_text)} characters",
                "Lines": len(extracted_text.split('\n'))
            }
        }
    except Exception as e:
        return {
            "type": "TXT",
            "title": uploaded_file.name,
            "summary": f"Error processing TXT: {str(e)}",
            "original_text": "",
            "meta": {}
        }

def transcribe_audio():
    recognizer = sr.Recognizer()
    try:
        with sr.Microphone() as source:
            st.info("Listening... Speak clearly (1-minute limit)")
            recognizer.adjust_for_ambient_noise(source)
            audio = recognizer.listen(source, timeout=60, phrase_time_limit=60)
            text = recognizer.recognize_google(audio)
            return text
    except sr.WaitTimeoutError:
        return "No speech detected"
    except Exception as e:
        return f"Error: {str(e)}"

def count_words(text):
    return len(text.split()) if text.strip() else 0

def get_badge_class(file_type):
    if file_type == "PDF":
        return "pdf-badge"
    elif file_type == "PPTX":
        return "pptx-badge"
    elif file_type == "TXT":
        return "txt-badge"
    return ""

# --- MAIN CONTENT ---
st.markdown("""
<div class="header-container">
    <h1 style="font-size: 2.5rem; color: #6e45e2; margin-bottom: 0.2rem; font-weight: 700;">Texer</h1>
    <p style="font-size: 1rem; color: #6e45e2; font-weight: 500; margin-bottom: 0.5rem;">
        AI-powered text summarizer that gets straight to the point
    </p>
</div>
""", unsafe_allow_html=True)

# Input Container
st.markdown('<div class="input-container">', unsafe_allow_html=True)
input_text = st.text_area(
    "", 
    height=250, 
    value=st.session_state.input_text, 
    key="text_input",
    label_visibility="collapsed", 
    placeholder="Paste your text here for instant summarization (or upload PDF/PPTX/TXT files)"
)

# Word count for input
input_word_count = count_words(input_text)
st.markdown(f'<div class="word-count">Words: {input_word_count}</div>', unsafe_allow_html=True)

# File uploader - allow multiple files
uploaded_files = st.file_uploader(
    "",  
    type=["pdf", "pptx", "txt"],
    help="Upload files for summarization (PDF/PPTX/TXT)",
    accept_multiple_files=True
)

# Process all uploaded files
if uploaded_files and not st.session_state.file_processed:
    st.session_state.file_summaries = []
    st.session_state.processing_files = True
    
    for uploaded_file in uploaded_files:
        try:
            if uploaded_file.type == "application/pdf":
                result = summarize_pdf(uploaded_file)
            elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                result = summarize_pptx(uploaded_file)
            elif uploaded_file.type == "text/plain":
                result = summarize_txt(uploaded_file)
            else:
                result = {
                    "type": "Unknown",
                    "title": uploaded_file.name,
                    "summary": "Unsupported file type",
                    "original_text": "",
                    "meta": {}
                }
            
            st.session_state.file_summaries.append(result)
            st.session_state.history.append({
                "title": result["title"],
                "summary": result["summary"],
                "original_text": result["original_text"],
                "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M"),
                "file_type": result["type"],
                "meta": result["meta"]
            })
        except Exception as e:
            st.error(f"Error processing {uploaded_file.name}: {str(e)}")
    
    st.session_state.file_processed = True
    st.session_state.processing_files = False
    st.rerun()

# Display summaries for uploaded files
if st.session_state.get("file_summaries"):
    st.markdown('<div class="divider"></div>', unsafe_allow_html=True)
    st.markdown('<div class="summary-section">', unsafe_allow_html=True)
    st.subheader("File Summaries")
    
    for idx, summary_data in enumerate(st.session_state.file_summaries):
        badge_class = get_badge_class(summary_data["type"])
        
        with st.expander(f"""
            <div>
                <span class="file-type-badge {badge_class}">{summary_data['type']}</span>
                {summary_data['title']}
            </div>
        """, expanded=True):
            st.markdown('<div class="file-summary">', unsafe_allow_html=True)
            
            # Display metadata
            st.markdown("**Document Details**")
            st.markdown(f"""
            <div class="file-meta">
                {''.join(f'<strong>{key}:</strong> {value}<br>' for key, value in summary_data['meta'].items())}
            </div>
            """, unsafe_allow_html=True)
            
            # Display summary
            st.markdown("**Summary**")
            st.markdown(summary_data["summary"])
            
            # Word count information
            original_words = count_words(summary_data["original_text"])
            summary_words = count_words(summary_data["summary"])
            reduction_pct = 100 - (summary_words / original_words * 100) if original_words > 0 else 0
            
            st.markdown(f"""
            <div class="word-count">
                Original: {original_words} words | 
                Summary: {summary_words} words | 
                Reduced by: {reduction_pct:.1f}%
            </div>
            """, unsafe_allow_html=True)
            
            # Copy button
            if st.button(f"Copy Summary", key=f"copy_{idx}"):
                st.session_state.copied_idx = idx
                st.rerun()
            
            if st.session_state.get("copied_idx") == idx:
                st.success("Summary copied to clipboard!")
                st.session_state.copied_idx = None
            
            st.markdown('</div>', unsafe_allow_html=True)
    
    st.markdown('</div>', unsafe_allow_html=True)

# Button container for side-by-side buttons
st.markdown('<div class="button-container">', unsafe_allow_html=True)
col1, col2 = st.columns(2)

with col1:
    if st.button("Generate Summary", 
                key="generate_btn",
                help="Create concise summary of your text",
                use_container_width=True):
        current_text = st.session_state.text_input if "text_input" in st.session_state else ""
        
        if not current_text.strip() and not st.session_state.input_text.strip():
            st.warning("Please enter some text or upload a file")
        else:
            text_to_summarize = current_text if current_text.strip() else st.session_state.input_text
            st.session_state.input_text = text_to_summarize
            
            with st.spinner("Analyzing content..."):
                summary = get_summary_from_gemini(
                    text_to_summarize,
                    st.session_state.summary_mode, 
                    st.session_state.length_factor
                )
                if not summary.startswith("Error"):
                    st.session_state.generated_summary = summary
                    st.session_state.show_summary = True
                    st.session_state.history.append({
                        "title": text_to_summarize[:50] + ("..." if len(text_to_summarize) > 50 else ""), 
                        "summary": summary, 
                        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M"),
                        "file_type": "TEXT",
                        "meta": {}
                    })
                    st.rerun()
                else:
                    st.error(summary)

with col2:
    if st.button("Voice Input", 
                key="record_btn",
                help="Dictate your text instead of typing",
                use_container_width=True):
        transcribed_text = transcribe_audio()
        if transcribed_text and not any(msg in transcribed_text for msg in ["Error", "No speech"]):
            st.session_state.input_text = transcribed_text
            st.session_state.file_processed = False
            st.rerun()
        elif transcribed_text:
            st.warning(transcribed_text)

st.markdown('</div>', unsafe_allow_html=True)  # Close button-container
st.markdown('</div>', unsafe_allow_html=True)  # Close input-container

# Text summary display
if st.session_state.show_summary and st.session_state.generated_summary:
    st.markdown('<div class="divider"></div>', unsafe_allow_html=True)
    st.markdown('<div class="summary-section">', unsafe_allow_html=True)
    st.subheader("Text Summary")
    st.markdown('<div class="summary-box">', unsafe_allow_html=True)
    
    st.markdown(st.session_state.generated_summary)
    
    # Word count for summary
    summary_word_count = count_words(st.session_state.generated_summary)
    st.markdown(f"""
    <div class="word-count">
        Original: {input_word_count} words | 
        Summary: {summary_word_count} words | 
        Reduced by: {100 - (summary_word_count/input_word_count*100 if input_word_count > 0 else 0):.1f}%
    </div>
    """, unsafe_allow_html=True)
    
    if st.button("Copy Summary", 
                key="copy_text_btn",
                help="Copy to clipboard",
                use_container_width=True):
        st.session_state.copied_text = True
        st.rerun()
    
    if st.session_state.get("copied_text", False):
        st.success("Text summary copied to clipboard!")
        st.session_state.copied_text = False
    
    st.markdown('</div>', unsafe_allow_html=True)
    st.markdown('</div>', unsafe_allow_html=True)

# --- SIDEBAR SETTINGS ---
with st.sidebar:
    # Settings section
    st.markdown('<div class="sidebar-section">', unsafe_allow_html=True)
    st.markdown('<div class="sidebar-title">Settings</div>', unsafe_allow_html=True)
    
    with st.expander("Customization", expanded=True):
        st.session_state.summary_mode = st.radio(
            "Summary Style:",
            ["Concise", "Brief", "Detailed"],
            help="Choose summary format"
        )
        st.session_state.length_factor = st.slider(
            "Summary Length:",
            0.1, 1.0, st.session_state.length_factor, 0.05,
            help="Adjust summary length (percentage of original)"
        )
    st.markdown('</div>', unsafe_allow_html=True)
    
    # History section
    st.markdown('<div class="sidebar-section">', unsafe_allow_html=True)
    st.markdown('<div class="sidebar-title">History</div>', unsafe_allow_html=True)
    
    with st.expander("Recent Summaries", expanded=True):
        if st.button("Clear History", 
                    key="clear_btn",
                    help="Delete all history",
                    use_container_width=True):
            st.session_state.history = []
            st.session_state.file_summaries = []
            st.session_state.show_summary = False
            st.rerun()
        
        if st.session_state.history:
            history_text = "\n\n".join(
                f"=== {item['title']} ===\n"
                f"Type: {item.get('file_type', 'TEXT')}\n"
                f"{''.join(f'{key}: {value}\n' for key, value in item.get('meta', {}).items())}"
                f"Summary: {item['summary']}\n\n"
                f"Date: {item['timestamp']}" 
                for item in st.session_state.history
            )
            st.download_button(
                label="Export History",
                data=history_text,
                file_name="texer_history.txt",
                mime="text/plain",
                use_container_width=True,
                help="Download summary history"
            )
            
            st.markdown("### Recent Items")
            for idx, item in enumerate(reversed(st.session_state.history[-5:])):
                badge_class = get_badge_class(item.get("file_type", "TEXT"))
                
                st.markdown('<div class="history-item">', unsafe_allow_html=True)
                st.markdown(f"""
                <div class="history-title">
                    <span class="file-type-badge {badge_class}">{item.get('file_type', 'TEXT')}</span>
                    {item["title"]}
                </div>
                """, unsafe_allow_html=True)
                st.markdown(f'<div class="history-date">{item["timestamp"]}</div>', unsafe_allow_html=True)
                
                if st.button("Load", key=f"load_{idx}", use_container_width=True):
                    if item.get("file_type") == "TEXT":
                        st.session_state.input_text = item.get('original_text', item['summary'])
                        st.session_state.generated_summary = item['summary']
                        st.session_state.show_summary = True
                    else:
                        # For files, we need to find the matching summary in file_summaries
                        for file_summary in st.session_state.file_summaries:
                            if file_summary["title"] == item["title"]:
                                st.session_state.file_summaries = [file_summary]
                                st.session_state.file_processed = True
                                break
                    st.rerun()
                st.markdown('</div>', unsafe_allow_html=True)
    st.markdown('</div>', unsafe_allow_html=True)

# Reset file processed flag after actions
if st.session_state.file_processed:
    st.session_state.file_processed = False
